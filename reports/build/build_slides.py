"""Build the presentation on the official HUST 2022 (red, 4:3) template.

Deletes the template's sample slides (the cover/closing artwork lives in the
layouts), then fills ~21 content slides through the branded layouts, embedding
the LaTeX formulas rendered by build_formulas.py and every analysis figure.
Run from repo root:

    python reports/build/build_slides.py

Output: reports/hust-house-price-slides.pptx
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent.parent
BRAND = ROOT / "reports" / "branding"
FIG = ROOT / "reports" / "figures"
FORM = BRAND / "formulas"

HUST_RED = RGBColor(0xC0, 0x00, 0x00)
INK = RGBColor(0x2B, 0x2B, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TITLE = "House Price Prediction for a Real Estate Listing Platform"
SUBTITLE = "From raw data to a monitored, deployed valuation tool"
COURSE = "IT5315E Business Analytics — Group Project"
TEAM = [
    ("Đỗ Hoàng Việt", "20252744M"),
    ("Nguyễn Minh Huyền", "20252007M"),
    ("Nguyễn Thị Huyền Trang", "20251322M"),
    ("Bùi Tá Cường", "20242430M"),
    ("Nguyễn Tiến Đức", "20252076M"),
]
APP_URL = "house-price-analytics-tjpb7ntxsd6kzbapp7bymmv.streamlit.app"


# --------------------------------------------------------------------------
# helpers
# --------------------------------------------------------------------------
def delete_slide(prs, slide):
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        if prs.part.related_part(sld.rId) == slide.part:
            prs.part.drop_rel(sld.rId)  # else the part is written twice
            xml_slides.remove(sld)


def set_title(slide, text, size=26):
    ph = slide.placeholders[0]
    ph.text = text
    for run in ph.text_frame.paragraphs[0].runs:
        run.font.size = Pt(size)
        run.font.bold = True


def bullets(tf, items, size=16, space_after=8):
    """Fill a text frame with (level, text, bold) bullet tuples."""
    tf.word_wrap = True
    for i, (level, text, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.space_after = Pt(space_after)
        for run in p.runs:
            run.font.size = Pt(size - 2 * level)
            run.font.bold = bold
            run.font.color.rgb = INK


def add_text(slide, left, top, width, height, lines, align=PP_ALIGN.LEFT):
    """Plain text box from (text, size, bold, color) tuples."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return box


def add_picture_fit(slide, path, left, top, max_w, max_h):
    """Insert a picture scaled to fit inside a box, centered in it."""
    w, h = Image.open(path).size
    scale = min(max_w / w, max_h / h)
    pw, ph = int(w * scale), int(h * scale)
    slide.shapes.add_picture(str(path), left + (max_w - pw) // 2,
                             top + (max_h - ph) // 2, pw, ph)


def add_table(slide, rows, left, top, width, col_widths=None, size=13):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top,
                                   width, Emu(1))
    table = shape.table
    if col_widths:
        total = sum(col_widths)
        for j, cw in enumerate(col_widths):
            table.columns[j].width = int(width * cw / total)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            if i == 0:  # header row in HUST red, not the Office default blue
                cell.fill.solid()
                cell.fill.fore_color.rgb = HUST_RED
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j and i else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(size)
                    run.font.bold = i == 0
                    run.font.color.rgb = WHITE if i == 0 else INK
    return table


def content_slide(prs, title):
    """Branded slide: red title bar + one content placeholder (layout 3)."""
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    set_title(slide, title)
    return slide


def two_content_slide(prs, title):
    """Branded slide with two side-by-side placeholders (layout 4)."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, title)
    return slide


def body_frame(slide, ph_idx):
    return slide.placeholders[ph_idx].text_frame


def figure_slide(prs, title, fig_path, notes, fig_box=None, note_size=14):
    """Full-width figure with takeaway bullets underneath (or beside)."""
    s = content_slide(prs, title)
    left, top, w, h = fig_box or (Inches(0.9), Inches(1.45), Inches(8.2), Inches(4.1))
    add_picture_fit(s, fig_path, left, top, w, h)
    box = s.shapes.add_textbox(Inches(0.7), top + h + Inches(0.15),
                               Inches(8.6), Inches(1.3))
    bullets(box.text_frame, notes, size=note_size, space_after=4)
    return s


def set_notes(slide, vi, en):
    """Attach a bilingual speaker note (Vietnamese first, then English)."""
    slide.notes_slide.notes_text_frame.text = f"VI: {vi}\n\nEN: {en}"


# Bilingual speaker scripts, one per slide, in slide-creation order.
NOTES = [
    # 1 — cover
    ("Chào thầy/cô và cả lớp. Nhóm em trình bày dự án dự đoán giá nhà cho một "
     "sàn bất động sản, đi trọn vòng đời từ dữ liệu thô đến mô hình được triển "
     "khai và giám sát. Nhóm gồm 5 thành viên.",
     "Good morning. We present a house-price prediction project for a real "
     "estate listing platform — the full lifecycle from raw data to a "
     "deployed, monitored model. Our team has five members."),
    # 2 — agenda
    ("Bài gồm 10 phần: bài toán kinh doanh, dữ liệu, phân tích khám phá, làm "
     "sạch, đặc trưng, so sánh mô hình, khoảng tin cậy, demo trực tiếp, giám "
     "sát, và khuyến nghị. Trọng tâm chấm điểm là độ chặt chẽ của từng bước.",
     "Ten parts: business problem, data, EDA, cleaning, features, model "
     "comparison, confidence intervals, live demo, monitoring, and "
     "recommendations. Grading emphasizes the rigor of each stage."),
    # 3 — business problem
    ("Định giá sai gây thiệt cả hai phía: hét giá cao đuổi người mua, để giá "
     "thấp mất tiền người bán. Môi giới cần ước tính đáng tin trước khi có "
     "thẩm định, và cần biết tin bao nhiêu. Cùng mức 180k, ±18k thì đăng bán "
     "tự tin, ±70k thì nên đi thẩm định — khoảng tin cậy đổi quyết định. KPI: "
     "sai số điển hình ≤10%, bất định được hiệu chỉnh, thiên lệch ≤5%, có "
     "triển khai chạy thật kèm giám sát.",
     "Mispricing hurts both sides: overpricing scares buyers, underpricing "
     "costs the seller. Agents need a trustworthy estimate before a formal "
     "appraisal — and need to know how much to trust it. At $180k, ±18k means "
     "list confidently, ±70k means appraise first — the range changes the "
     "decision. KPIs: typical error ≤10%, calibrated uncertainty, bias ≤5%, "
     "a live deployment with monitoring."),
    # 4 — CRISP-DM + KPI bridge
    ("Nhóm theo phương pháp luận CRISP-DM: hiểu bài toán kinh doanh trước, rồi "
     "mọi metric kỹ thuật mới sinh ra từ một KPI. Đọc bảng từ trên xuống: cần "
     "ước tính đủ tốt → MAPE; không làm lệch thống kê thị trường → thiên lệch; "
     "nói cho người dùng độ tin cậy → độ phủ; đúng cỡ bất định → độ rộng "
     "khoảng; bắt sai số lớn và giải thích giá → RMSE và R². Nói cách khác, "
     "RMSE/MAE/MAPE/R² không báo cáo cho có — mỗi cái trả lời một mục tiêu.",
     "We follow CRISP-DM: business understanding first, then every technical "
     "metric derives from a KPI. Read the table top-down: a good-enough "
     "estimate → MAPE; don't distort market stats → bias; tell the user how "
     "much to trust it → coverage; right-size uncertainty → interval width; "
     "catch big misses and explain price → RMSE and R². So RMSE/MAE/MAPE/R² "
     "aren't reported for their own sake — each answers a goal."),
    # 5 — data strategy
    ("Nền là bộ Kaggle Ames: 1.460 giao dịch, 80 thuộc tính. Nhóm mở rộng "
     "thêm 12 trường ngữ cảnh tổng hợp, mỗi trường có tài liệu đầy đủ. Quan "
     "trọng: trục thời gian neo vào ngày bán thật, 2006 đến 7/2010 — nên có "
     "cả giai đoạn khủng hoảng. Toàn bộ 175 giao dịch 2010 được giữ lại làm "
     "luồng dữ liệu đến cho phần giám sát.",
     "Base is Kaggle Ames: 1,460 sales, 80 attributes. We add 12 synthetic "
     "contextual fields, each fully documented. Key point: the time axis is "
     "anchored to the real sale dates, 2006 to July 2010 — so it spans the "
     "crisis. All 175 sales of 2010 are held out as the incoming stream for "
     "monitoring."),
    # 6 — synthetic credibility
    ("Dữ liệu tổng hợp không sinh ngẫu nhiên vô căn cứ mà theo điều kiện có "
     "logic: khu đắt tiền gần trường hơn (tương quan −0.60), còn giao thông "
     "cố ý không định giá. Mọi giả định đều được test tự động kiểm chứng và "
     "sinh lại nguyên vẹn từ seed cố định. Đợt phục hồi +9.5% chỉ nằm trong "
     "luồng giám sát — đó là tín hiệu drift được thiết kế có chủ ý.",
     "The synthetic data isn't random — it's generated conditionally with "
     "logic: pricier areas sit closer to schools (r = −0.60), while transit "
     "is deliberately unpriced. Every assumption is checked by automated "
     "tests and reproduces exactly from a fixed seed. The +9.5% rebound "
     "lives only in the monitoring stream — a deliberately designed drift "
     "signal."),
    # 7 — designed imperfections
    ("Để phần làm sạch là việc thật chứ không diễn, nhóm cố ý tiêm 5 loại lỗi "
     "vào các trường tổng hợp: thiếu giá trị, nhãn tự do, giá trị sentinel, "
     "sai đơn vị, và trùng dòng. Chỉ tiêm vào trước 2010 nên luồng giám sát "
     "vẫn sạch — drift không bị lẫn với bẩn. Vì biết trước sự thật nền, chất "
     "lượng làm sạch chứng minh được.",
     "So cleaning is real work, we deliberately inject five kinds of dirt "
     "into the synthetic fields: missing values, free-text labels, sentinel "
     "values, wrong units, and duplicate rows. Injected only before 2010, so "
     "the monitoring stream stays clean — drift is never confused with dirt. "
     "Because we know the ground truth, cleaning quality is provable."),
    # 8 — EDA target
    ("Giá bán lệch phải mạnh: trung vị 163k nhưng đuôi kéo tới 779k, độ lệch "
     "1.93. Sau khi lấy log, độ lệch còn 0.12 — gần chuẩn. Nên toàn bộ mô "
     "hình dự đoán trên thang log rồi đổi ngược về đô-la.",
     "Sale price is strongly right-skewed: median $163k but a tail to $779k, "
     "skew 1.93. After a log transform, skew drops to 0.12 — near-normal. So "
     "all models predict on the log scale and back-transform to dollars."),
    # 9 — EDA drivers
    ("Hai yếu tố mạnh nhất là chất lượng tổng thể (0.79) và diện tích ở "
     "(0.71) — giá trung vị tăng gấp bốn từ chất lượng 3 lên 9. Các tiện ích "
     "tổng hợp hành xử đúng thiết kế: gần trường và bệnh viện làm tăng giá, "
     "giao thông thì không. Chi phí cải tạo tương quan 0.69 nhưng một phần là "
     "dẫn xuất, nên sẽ kiểm ở bước VIF.",
     "The two strongest drivers are overall quality (0.79) and living area "
     "(0.71) — median price quadruples from quality 3 to 9. Synthetic "
     "amenities behave as designed: school and hospital proximity raise "
     "price, transit doesn't. Renovation cost correlates 0.69 but is partly "
     "derivative, so we check it at the VIF step."),
    # 10 — EDA outliers
    ("Hai căn trên 4.000 sqft bán thấp hơn hẳn xu hướng — đây là các partial "
     "sale nổi tiếng của bộ Ames, tức giao dịch không theo thị trường. Nhóm "
     "loại chúng ở bước làm sạch kèm lý do rõ ràng, không xoá tuỳ tiện.",
     "Two homes over 4,000 sq ft sold far below trend — the well-known Ames "
     "partial sales, i.e. non-market transactions. We drop them in cleaning "
     "with an explicit justification, not arbitrarily."),
    # 11 — EDA location
    ("Vị trí là đòn bẩy giá gấp 3,5 lần: từ 90k ở khu rẻ nhất đến 313k ở khu "
     "đắt nhất. Phần lớn chênh lệch giải thích được bằng khoảng cách tiện "
     "ích. Điểm actionable cho sàn: gần trường được định giá, còn gần bến xe "
     "thì không.",
     "Location is a 3.5× price lever: $90k in the cheapest neighborhood to "
     "$313k in the priciest. Most of the gap is explained by amenity "
     "distances. Actionable insight for the platform: proximity to schools "
     "is priced in, bus access is not."),
    # 12 — multicollinearity
    ("Có bốn cặp biến tương quan trên 0.8 — chúng nói cùng một điều. Nhóm giữ "
     "lại một biến cho mỗi khái niệm rồi kiểm bằng hệ số phóng đại phương sai "
     "VIF. VIF lớn nhất chỉ 3.9, dưới ngưỡng 10 — nên các mô hình tuyến tính "
     "ổn định và hệ số diễn giải được.",
     "Four variable pairs correlate above 0.8 — they say the same thing. We "
     "keep one variable per concept, then verify with variance inflation "
     "factors. Max VIF is only 3.9, below the threshold of 10 — so the "
     "linear models are stable and their coefficients interpretable."),
    # 13 — cleaning evidence
    ("Bảng này in ra mỗi lần chạy pipeline: mỗi hàng là một lỗi được chẩn "
     "đoán nguyên nhân và cách xử lý, kèm số lượng trước/sau. Ví dụ 12 tin "
     "trùng giữ bản sớm nhất, 51 nhãn cải tạo chuẩn hoá, sai đơn vị mét chia "
     "1000. Vào 1.472 dòng, ra 1.458 — mọi sửa đổi đều truy vết được.",
     "This table prints on every pipeline run: each row is a diagnosed "
     "problem, its fix, and before/after counts. E.g. 12 duplicate listings "
     "keep the earliest, 51 renovation labels normalized, metre units "
     "divided by 1000. 1,472 rows in, 1,458 out — every change is "
     "traceable."),
    # 14 — features & leakage
    ("Nhóm tạo các đặc trưng dẫn xuất như tuổi nhà, tổng số phòng tắm, điểm "
     "tiện ích. Kỷ luật quan trọng nhất là chống rò rỉ: chỉ dùng thứ môi giới "
     "biết tại thời điểm định giá. Loại days-on-market vì là kết quả sau khi "
     "đăng, và loại chỉ số giá thị trường vì công bố trễ — đưa vào sẽ nhìn "
     "trộm đợt phục hồi 2010. Hợp đồng đặc trưng cuối: 16 biến số + 6 biến "
     "phân loại.",
     "We build derived features like property age, total baths, amenity "
     "score. The key discipline is anti-leakage: only what an agent knows at "
     "valuation time. We exclude days-on-market — an outcome of listing — "
     "and the market price index — published with a lag; feeding it would "
     "let the model peek at the 2010 rebound. Final feature contract: 16 "
     "numeric + 6 categorical."),
    # 15 — model comparison
    ("Năm mô hình so sánh bằng CV 5-fold. Ridge thắng về điểm dự đoán: MAPE "
     "8.99%, R² 0.916. Cây kém hơn ở đây vì log-giá gần tuyến tính với các "
     "yếu tố và cỡ mẫu ~1.300 ưu tiên mô hình phương sai thấp. Nhưng — như "
     "slide sau nói — độ chính xác điểm chưa phải sản phẩm cuối.",
     "Five models compared by 5-fold CV. Ridge wins on point accuracy: MAPE "
     "8.99%, R² 0.916. Trees do worse here because log-price is near-linear "
     "in the drivers and n ≈ 1,300 favors low-variance models. But — as the "
     "next slide argues — point accuracy isn't the final deliverable."),
    # 16 — residuals
    ("Phân tích phần dư: không có xu hướng theo giá trị dự đoán, các decile "
     "2–9 không thiên lệch — mô hình lành mạnh ở phần lớn dải giá. Điểm yếu "
     "thành thật: decile rẻ nhất bị dự đoán cao +13%, decile đắt nhất hơi "
     "thấp −4.7% — hiện tượng co ngót cổ điển. Nhóm ghi nhận là hạn chế và "
     "khoảng tin cậy được nới rộng đúng ở những vùng đó.",
     "Residual analysis: no trend versus fitted values, deciles 2–9 unbiased "
     "— the model is healthy across most of the range. Honest weakness: the "
     "cheapest decile is over-predicted by +13%, the priciest slightly under "
     "by −4.7% — classic shrinkage. We flag it as a limitation, and "
     "intervals widen exactly there."),
    # 17 — model choice
    ("Đây là lập luận trung tâm: bài toán yêu cầu một khoảng tin cậy cho từng "
     "căn, không chỉ một con số. Nên nhóm triển khai ba mô hình LightGBM cho "
     "phân vị 10/50/90 huấn luyện bằng pinball loss — cho khoảng phụ thuộc "
     "đặc trưng một cách tự nhiên. Đây là một họ mô hình mạch lạc, không chắp "
     "điểm Ridge vào khoảng của mô hình khác. Cái giá 9.78% so với 8.99% "
     "MAPE là không đáng kể so với KPI biên đàm phán, và nhóm báo cáo công "
     "khai chứ không giấu.",
     "This is the central argument: the task requires a per-home confidence "
     "range, not just a number. So we deploy three LightGBM models for the "
     "10/50/90 quantiles trained on the pinball loss — feature-dependent "
     "intervals natively. It's one coherent model family, not a Ridge point "
     "stitched onto someone else's band. The cost, 9.78% vs 8.99% MAPE, is "
     "immaterial against the negotiation-margin KPI — and we report it "
     "openly."),
    # 18 — conformal calibration
    ("Đây là phần thành thật nhất. Khoảng thô ban đầu tuyên bố phủ 80% nhưng "
     "thực tế chỉ phủ 59% trên dữ liệu chưa thấy — tức nó nói dối. Nhóm dùng "
     "conformalized quantile regression để đo đúng độ thiếu hụt rồi nới "
     "khoảng theo phân vị thực nghiệm của điểm số. Sau hiệu chỉnh, độ phủ đo "
     "được 80.5%, khớp mục tiêu. Độ rộng co giãn hợp lý: ~35k cho nhà điển "
     "hình, ~78k cho nhà cao cấp — khoảng rộng nghĩa là nên đi thẩm định.",
     "This is our most honest part. The raw band claimed 80% coverage but "
     "actually covered only 59% on unseen data — it lied. We use "
     "conformalized quantile regression to measure the shortfall, then widen "
     "the band by the empirical quantile of the scores. After calibration, "
     "measured coverage is 80.5%, on target. Width scales sensibly: ~$35k "
     "for typical homes, ~$78k for high-value ones — a wide range means "
     "'consider an appraisal'."),
    # 19 — live demo
    ("Giờ em demo công cụ chạy thật trên Streamlit. Kịch bản A — nhà điển "
     "hình ở NAmes cho khoảng hẹp. Kịch bản B — biệt thự ở NridgHt cho khoảng "
     "rộng kèm khuyến nghị thẩm định. Cùng mô hình đó cũng chạy sau một REST "
     "API FastAPI. Để kiểm soát rủi ro, em đã làm nóng link trước và có bản "
     "chạy local dự phòng.",
     "Now a live demo of the tool on Streamlit. Scenario A — a typical home "
     "in NAmes gives a tight range. Scenario B — a mansion in NridgHt gives "
     "a wide range plus appraisal advice. The same model also runs behind a "
     "FastAPI REST endpoint. For risk control, I warmed the link beforehand "
     "and have an identical local fallback."),
    # 20 — monitoring design
    ("Mô hình tự theo dõi bản thân qua bốn trigger, tất cả định nghĩa trước "
     "khi replay dữ liệu để không bị chỉnh cho khớp: T1 drift dữ liệu, T2 "
     "hiệu năng RMSE, T3 thiên lệch hệ thống, T4 sức khoẻ khoảng tin cậy. "
     "Dùng cửa sổ trượt 3 tháng vì từng tháng riêng lẻ quá ít mẫu — cửa sổ 10 "
     "dòng gắn cờ mọi biến là drift, đã kiểm chứng là nhiễu.",
     "The model watches itself through four triggers, all defined before "
     "replaying the data so they're not tuned to fit: T1 data drift, T2 "
     "performance RMSE, T3 systematic bias, T4 interval health. We use "
     "3-month rolling windows because single months are too small — a 10-row "
     "window flags every feature as drifted, verified to be noise."),
    # 21 — 2010 replay
    ("Cho luồng 2010 chạy qua từng tháng: tháng 1–2 lành mạnh. Tháng 3 RMSE "
     "bật lên và drift vượt ngưỡng — tín hiệu retrain đầu tiên. Từ tháng 4 "
     "đến 7, thiên lệch trôi từ −0.4% xuống −7.9% khi thị trường phục hồi "
     "nhanh hơn cửa sổ huấn luyện, độ phủ tụt còn ~73%. Kết luận: retrain từ "
     "tháng 3/2010. Bài học: drift đầu vào nói có gì đó đổi, nhưng chỉ giám "
     "sát sai số và thiên lệch mới chứng minh điều đó quan trọng.",
     "Running the 2010 stream month by month: January–February healthy. "
     "March, RMSE spikes and drift crosses the threshold — the first retrain "
     "signal. From April to July, bias drifts from −0.4% to −7.9% as the "
     "market rebounds faster than the training window, and coverage sags to "
     "~73%. Verdict: retrain from March 2010. Lesson: input drift says "
     "something changed, but only error and bias monitoring prove it "
     "matters."),
    # 22 — recommendations & limitations
    ("Khuyến nghị cho sàn: định giá trước hết dựa vào chất lượng và vị trí; "
     "luôn hiển thị khoảng tin cậy vì nó biến hộp đen thành công cụ đàm phán; "
     "coi mô hình là hàng dễ hỏng và lên ngân sách retrain hàng quý; đưa nhà "
     "có khoảng rộng sang dịch vụ thẩm định. Hạn chế: lớp ngữ cảnh là tổng "
     "hợp, giá ở mức Ames 2010, decile rẻ bị dự đoán cao, và retrain hiện do "
     "trigger đề xuất chứ chưa tự động.",
     "Recommendations for the platform: price on quality and location first; "
     "always show the range because it turns a black box into a negotiation "
     "aid; treat the model as perishable and budget quarterly retraining; "
     "route wide-interval homes to appraisal. Limitations: the contextual "
     "layer is synthetic, prices are at the 2010 Ames level, the cheap "
     "decile is over-predicted, and retraining is trigger-recommended but "
     "not yet automated."),
    # 23 — closing
    ("Cảm ơn thầy/cô và cả lớp đã lắng nghe. Công cụ chạy thật và mã nguồn ở "
     "hai link này. Nhóm em sẵn sàng trả lời câu hỏi.",
     "Thank you for listening. The live tool and source code are at these "
     "two links. We're happy to take questions."),
]


# --------------------------------------------------------------------------
# deck
# --------------------------------------------------------------------------
def build():
    prs = Presentation(BRAND / "hust-template-4x3.pptx")
    for slide in list(prs.slides):
        delete_slide(prs, slide)

    # ---- 1. cover (layout 2 = white title-page artwork) --------------------
    cover = prs.slides.add_slide(prs.slide_layouts[2])
    add_text(cover, Inches(0.55), Inches(1.25), Inches(6.4), Inches(1.9), [
        (TITLE, 28, True, HUST_RED)])
    add_text(cover, Inches(0.55), Inches(3.15), Inches(6.2), Inches(0.8), [
        (SUBTITLE, 15, True, INK),
        (COURSE, 13, False, INK)])
    team_lines = [(f"{name}  —  {sid}", 12, False, INK) for name, sid in TEAM]
    add_text(cover, Inches(0.55), Inches(4.15), Inches(6.0), Inches(2.2),
             team_lines)

    # ---- 2. agenda ----------------------------------------------------------
    s = two_content_slide(prs, "Agenda")
    bullets(body_frame(s, 1), [
        (0, "1. The business problem", False),
        (0, "2. Data: real base + synthetic context", False),
        (0, "3. Exploratory analysis — what drives price", False),
        (0, "4. Data cleaning with evidence", False),
        (0, "5. Features & multicollinearity", False),
    ], size=16)
    bullets(body_frame(s, 2), [
        (0, "6. Model comparison (5-fold CV)", False),
        (0, "7. Calibrated confidence intervals", False),
        (0, "8. Live demo — the valuation tool", False),
        (0, "9. Monitoring & retraining triggers", False),
        (0, "10. Recommendations & limitations", False),
    ], size=16)

    # ---- 3. business problem ------------------------------------------------
    s = content_slide(prs, "The business problem")
    bullets(body_frame(s, 13), [
        (0, "Mispricing hurts both sides: inflated asks scare buyers away; "
            "underpricing leaves the seller's money on the table.", False),
        (0, "Agents need a defensible estimate long before a formal appraisal "
            "— and need to know how much to trust it.", False),
        (0, "Business question: given a property's characteristics and "
            "location, what is its market price — and how confident is that "
            "estimate?", True),
        (0, "$180k ± $18k → list with confidence.  $180k ± $70k → order an "
            "appraisal first. The range changes the decision.", False),
        (0, "KPIs: typical error ≤ ~10% (negotiation margin); calibrated "
            "uncertainty; systematic bias ≤ 5%; a working deployment with "
            "monitoring.", False),
    ], size=15)

    # ---- 3b. CRISP-DM + KPI-to-metric bridge --------------------------------
    s = content_slide(prs, "From business goals to metrics (CRISP-DM)")
    add_text(s, Inches(0.6), Inches(1.35), Inches(8.8), Inches(0.9), [
        ("We follow CRISP-DM: business understanding comes first. Every "
         "technical metric below exists to answer a business KPI — RMSE / MAE "
         "/ MAPE / R² are not reported for their own sake.", 14, False, INK),
    ])
    add_table(s, [
        ["Business goal / KPI", "Technical metric", "Target · result"],
        ["Estimate agents can act on", "MAPE (with MAE in $)", "≤ ~10%  ·  9.78%"],
        ["Don't skew market statistics", "Signed bias, % of median", "≤ 5%  ·  monitored"],
        ["Tell user how much to trust it", "Interval coverage (unseen)", "≈ 80%  ·  80.5%"],
        ["Right-size uncertainty per home", "Interval width by segment", "$35k vs $78k"],
        ["Catch big misses / explain price", "RMSE  ·  R²", "$24.4k  ·  0.88–0.92"],
    ], Inches(0.5), Inches(2.45), Inches(9.0),
       col_widths=[3.6, 3.0, 2.4], size=13)
    box = s.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(8.8), Inches(0.7))
    bullets(box.text_frame, [
        (0, "Read top-down: the deliverable is a calibrated, unbiased, "
            "decision-changing estimate — the metrics just measure it.", True)],
        size=13)

    # ---- 4. data strategy ----------------------------------------------------
    s = two_content_slide(prs, "Data: real base + documented synthetic context")
    bullets(body_frame(s, 1), [
        (0, "Base: Kaggle Ames — 1,460 sales, 80 attributes.", False),
        (0, "Extension: 12 synthetic contextual fields, each documented "
            "(type, unit, range, assumption).", False),
        (0, "Time axis anchored to the data's real sale dates: "
            "Jan 2006 – Jul 2010.", True),
        (0, "Macro overlay: interest-rate series; price index with boom, "
            "−10% crisis, +9.5% rebound in 2010.", False),
        (0, "All 175 sales of 2010 held out as the incoming stream for "
            "monitoring.", False),
    ], size=15)
    add_picture_fit(s, FIG / "02_market_over_time.png",
                    Inches(5.1), Inches(1.5), Inches(4.6), Inches(5.4))

    # ---- 5. synthetic credibility ---------------------------------------------
    s = content_slide(prs, "Why the synthetic data can be trusted")
    bullets(body_frame(s, 13), [
        (0, "Behavioral assumptions, generated conditionally:", True),
        (1, "Pricier neighborhoods sit closer to schools (realized r = "
            "−0.60); transit access deliberately unpriced (r ≈ 0.1).", False),
        (1, "Renovation costs consistent with the real remodel-year column; "
            "days-on-market lengthens in the 2008–09 cold market.", False),
        (0, "Prices re-expressed under the synthetic market index "
            "(original Kaggle price preserved for audit):", True),
    ], size=15)
    add_picture_fit(s, FORM / "price_index.png",
                    Inches(0.9), Inches(3.85), Inches(8.2), Inches(0.85))
    box = s.shapes.add_textbox(Inches(0.65), Inches(4.9), Inches(8.7), Inches(1.7))
    bullets(box.text_frame, [
        (0, "Every documented claim is enforced by automated tests; "
            "generation reproduces byte-for-byte from a fixed seed.", False),
        (0, "The +9.5% rebound lives only inside the monitoring stream — the "
            "designed, explainable drift signal.", True),
    ], size=15, space_after=4)

    # ---- 6. injected imperfections ---------------------------------------------
    s = content_slide(prs, "Designed imperfections — so cleaning is real work")
    add_table(s, [
        ["Field", "Injected problem", "Rate"],
        ["dist_school_km", "Missing — blank at listing entry", "2% (26 rows)"],
        ["renovated", "Free-text labels: Y / yes / YES / N / no / NO", "4% (51)"],
        ["days_on_market", 'Legacy sentinel 999 = "unknown"', "0.4% (5)"],
        ["dist_transit_km", "Entered in metres instead of km", "0.3% (4)"],
        ["whole row", "Home re-posted under a new listing id", "12 rows"],
    ], Inches(0.7), Inches(1.6), Inches(8.6), col_widths=[2.2, 5.2, 1.6], size=13)
    box = s.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(8.6), Inches(1.7))
    bullets(box.text_frame, [
        (0, "Injected only into synthetic fields, only before 2010 — the "
            "monitoring stream stays clean, so drift is never confused with "
            "dirt.", False),
        (0, "Ground truth is known → cleaning quality can be proven, not "
            "assumed.", True),
    ], size=14, space_after=4)

    # ---- 7. EDA: target -----------------------------------------------------------
    s = figure_slide(prs, "EDA 1/4 — the target is right-skewed",
                     FIG / "01_target_distribution.png", [
        (0, "Median $162.9k vs mean $181.4k, tail to $779k; skew 1.93 → 0.12 "
            "after log. All models predict:", False)],
        fig_box=(Inches(1.0), Inches(1.5), Inches(8.0), Inches(3.3)))
    add_picture_fit(s, FORM / "target.png",
                    Inches(2.9), Inches(5.75), Inches(4.2), Inches(0.8))

    # ---- 8. EDA: drivers ------------------------------------------------------------
    s = two_content_slide(prs, "EDA 2/4 — what drives price")
    add_picture_fit(s, FIG / "03_top_price_drivers.png",
                    Inches(0.3), Inches(1.6), Inches(4.9), Inches(4.4))
    bullets(body_frame(s, 2), [
        (0, "Quality (r = 0.79) and living area (0.71) dominate — median "
            "price quadruples from quality 3 ($85k) to 9 ($346k).", False),
        (0, "Synthetic amenities behave as designed: school −0.60, hospital "
            "−0.58, transit ≈ +0.1.", False),
        (0, "renovation_cost (0.69) is partly derivative — generated from "
            "area × quality; flagged for the VIF check.", False),
    ], size=14)

    # ---- 9. EDA: outliers -----------------------------------------------------------
    figure_slide(prs, "EDA 3/4 — quality, size, and two notorious outliers",
                 FIG / "04_quality_area_outliers.png", [
        (0, "Two homes > 4,000 sq ft sold far below trend (Ids 524, 1299) — "
            "known Ames partial sales, dropped in cleaning with "
            "justification.", False)])

    # ---- 10. EDA: location ------------------------------------------------------------
    s = content_slide(prs, "EDA 4/4 — location is a 3.5× price lever")
    add_picture_fit(s, FIG / "05_neighborhood_prices.png",
                    Inches(0.5), Inches(1.45), Inches(9.0), Inches(2.5))
    add_picture_fit(s, FIG / "06_amenity_distances.png",
                    Inches(0.5), Inches(4.05), Inches(9.0), Inches(2.1))
    box = s.shapes.add_textbox(Inches(0.7), Inches(6.25), Inches(8.6), Inches(0.7))
    bullets(box.text_frame, [
        (0, "Medians $90k (MeadowV) → $313k (NridgHt). School proximity is "
            "priced in; bus access is not — an actionable platform insight.",
         True)], size=13)

    # ---- 11. multicollinearity ----------------------------------------------------------
    s = two_content_slide(prs, "Redundant variables — mapped, then pruned")
    add_picture_fit(s, FIG / "07_correlation_heatmap.png",
                    Inches(0.25), Inches(1.5), Inches(5.2), Inches(5.2))
    bullets(body_frame(s, 2), [
        (0, "Four pairs above |r| = 0.8: garage cars/area (0.88), rooms/area "
            "(0.83), 1st-floor/basement (0.82), build/garage year (0.83).",
         False),
        (0, "Kept one variable per concept; verified with variance "
            "inflation:", False),
    ], size=14)
    add_picture_fit(s, FORM / "vif.png",
                    Inches(6.1), Inches(3.7), Inches(2.6), Inches(1.0))
    box = s.shapes.add_textbox(Inches(5.5), Inches(4.9), Inches(4.1), Inches(1.4))
    bullets(box.text_frame, [
        (0, "Max VIF = 3.9 across 16 numeric features (threshold 10) — "
            "stable, interpretable linear baselines.", True)], size=14)

    # ---- 12. cleaning ---------------------------------------------------------------------
    s = content_slide(prs, "Cleaning with before / after evidence")
    add_table(s, [
        ["Problem (diagnosed cause)", "Before", "After"],
        ["Duplicate listings (re-posted homes) — keep earliest", "12", "0"],
        ["Free-text renovation labels (Y/yes/NO/…)", "51", "0"],
        ["Days-on-market sentinel 999 → quarter-median impute", "5", "0"],
        ["Transit distance entered in metres → ÷1000", "4", "0"],
        ["Missing school distance → neighborhood median", "26", "0"],
        ['Kaggle "no such feature" NAs → explicit category', "7,480", "0"],
        ["LotFrontage → neighborhood median", "259", "0"],
        ["Partial-sale outliers (Ids 524, 1299) → dropped", "2 rows", "—"],
    ], Inches(0.7), Inches(1.55), Inches(8.6), col_widths=[6, 1.2, 1.2], size=12)
    box = s.shapes.add_textbox(Inches(0.7), Inches(5.95), Inches(8.6), Inches(0.9))
    bullets(box.text_frame, [
        (0, "1,472 listings in → 1,458 out. Every fix targets a diagnosed "
            "cause; the pipeline prints this table on every run.", False),
    ], size=13)

    # ---- 13. features & leakage -----------------------------------------------------------
    s = content_slide(prs, "Feature engineering & leakage discipline")
    bullets(body_frame(s, 13), [
        (0, "Derived: property age, years since remodel, total baths, "
            "renovation flag & cost, amenity score, listing month.", False),
        (0, "Only what an agent knows at valuation time:", True),
        (1, "Excluded: days-on-market — an outcome of the listing, unknown "
            "on pricing day.", False),
        (1, "Excluded: market price index — published with a lag; feeding "
            "the current value would be look-ahead leakage (and would let "
            "the model “see” the 2010 rebound).", False),
        (1, "Included: interest rate — public on listing day.", False),
        (0, "Final feature contract: 16 numeric + 6 categorical variables.",
         True),
    ], size=15)

    # ---- 14. model comparison ----------------------------------------------------------------
    s = content_slide(prs, "Model comparison — 5-fold cross-validation")
    add_table(s, [
        ["Model", "RMSE", "MAE", "MAPE", "R²"],
        ["Linear Regression", "$23,409", "$15,806", "9.11%", "0.914"],
        ["Ridge (best point model)", "$23,195", "$15,573", "8.99%", "0.916"],
        ["Lasso", "$23,324", "$15,670", "9.04%", "0.915"],
        ["Random Forest", "$29,509", "$18,469", "10.43%", "0.863"],
        ["LightGBM (deployed, quantile)", "$27,441", "$17,334", "9.78%", "0.882"],
    ], Inches(0.8), Inches(1.5), Inches(8.4), col_widths=[3.2, 1.2, 1.2, 1, 1], size=13)
    add_picture_fit(s, FORM / "metrics_errors.png",
                    Inches(1.3), Inches(4.25), Inches(7.4), Inches(0.95))
    add_picture_fit(s, FORM / "metrics_relative.png",
                    Inches(1.3), Inches(5.25), Inches(7.4), Inches(0.95))
    box = s.shapes.add_textbox(Inches(0.7), Inches(6.35), Inches(8.6), Inches(0.7))
    bullets(box.text_frame, [
        (0, "Log-price is near-linear in the drivers and n ≈ 1,300 favors "
            "low-variance models — Ridge wins the point contest.", False)],
        size=12)

    # ---- 15. residuals --------------------------------------------------------------------------
    figure_slide(prs, "Residual analysis — where the model is honest, and not",
                 FIG / "08_residual_analysis.png", [
        (0, "No trend vs fitted values; deciles 2–9 unbiased.", False),
        (0, "Cheapest decile over-predicted (+13%), priciest slightly under "
            "(−4.7%) — classic shrinkage; flagged as a limitation and "
            "covered by wider intervals exactly there.", False)],
        fig_box=(Inches(0.7), Inches(1.5), Inches(8.6), Inches(3.6)),
        note_size=13)

    # ---- 16. why quantile model ships ------------------------------------------------------------
    s = content_slide(prs, "Model choice: accuracy is not the deliverable")
    bullets(body_frame(s, 13), [
        (0, "The business asks for a per-home confidence range, not just a "
            "number.", True),
        (0, "Deployed: three LightGBM models (p10 / p50 / p90) trained on "
            "the pinball loss — feature-dependent intervals natively:", False),
    ], size=15)
    add_picture_fit(s, FORM / "pinball.png",
                    Inches(0.8), Inches(3.0), Inches(8.4), Inches(0.85))
    box = s.shapes.add_textbox(Inches(0.65), Inches(4.1), Inches(8.7), Inches(2.3))
    bullets(box.text_frame, [
        (0, "One coherent model family — no stitching a Ridge point onto "
            "someone else's interval (mismatched band, two failure modes).",
         False),
        (0, "Cost: 9.78% vs 8.99% MAPE — immaterial against the negotiation-"
            "margin KPI; the comparison is reported, not hidden.", False),
    ], size=15, space_after=6)

    # ---- 17. conformal calibration ------------------------------------------------------------------
    s = content_slide(prs, "The naive band lies — conformal calibration fixes it")
    box = s.shapes.add_textbox(Inches(0.65), Inches(1.45), Inches(8.7), Inches(0.8))
    bullets(box.text_frame, [
        (0, "On 257 held-out sales, the raw p10–p90 band claimed 80% but "
            "covered only 59.1%. CQR measures the shortfall…", True)], size=15)
    add_picture_fit(s, FORM / "cqr_score.png",
                    Inches(1.3), Inches(2.45), Inches(7.4), Inches(0.8))
    box = s.shapes.add_textbox(Inches(0.65), Inches(3.45), Inches(8.7), Inches(0.55))
    bullets(box.text_frame, [
        (0, "…and widens the band by the empirical quantile of the scores:",
         False)], size=15)
    add_picture_fit(s, FORM / "cqr_interval.png",
                    Inches(0.9), Inches(4.05), Inches(8.2), Inches(0.85))
    box = s.shapes.add_textbox(Inches(0.65), Inches(5.1), Inches(8.7), Inches(1.6))
    bullets(box.text_frame, [
        (0, "Verified coverage after calibration: 80.5% (target 80%).", True),
        (0, "Width scales honestly: ≈$35k typical · ≈$41k mid · ≈$78k "
            "high-value tercile — wide range ⇒ “consider an appraisal”.",
         False),
    ], size=15, space_after=6)

    # ---- 18. live demo ----------------------------------------------------------------------------------
    s = two_content_slide(prs, "Live demo — the valuation tool")
    bullets(body_frame(s, 1), [
        (0, "Open the live tool", True),
        (0, "Scenario A — typical: NAmes, 1,500 sq ft, quality 5 → tight "
            "range (±13%).", False),
        (0, "Scenario B — high-value: NridgHt, 3,200 sq ft, quality 9 → "
            "wide range + appraisal advice.", False),
        (0, "Same model behind a REST API: POST /predict (FastAPI, OpenAPI "
            "docs).", False),
        (0, "Risk control: link warmed before the talk; identical local "
            "fallback ready.", False),
    ], size=14)
    # make the first bullet an actual clickable hyperlink
    demo_run = body_frame(s, 1).paragraphs[0].runs[0]
    demo_run.hyperlink.address = f"https://{APP_URL}/"
    demo_run.font.color.rgb = HUST_RED
    add_picture_fit(s, BRAND / "app_screenshot.png",
                    Inches(4.9), Inches(1.6), Inches(4.85), Inches(5.2))

    # ---- 19. monitoring design ----------------------------------------------------------------------------
    s = content_slide(prs, "Monitoring design — four triggers, fixed in advance")
    add_table(s, [
        ["Trigger", "Definition", "Why it matters"],
        ["T1  Data drift", "> 30% of features drift (Evidently, ≥40-row window)",
         "The world changed"],
        ["T2  Performance", "Rolling RMSE > 1.25 × baseline ($24.4k)",
         "Errors grew materially"],
        ["T3  Systematic bias", "|mean error| > 5% of median price",
         "One-directional mispricing distorts the platform"],
        ["T4  Interval health", "80% range covers < 65%",
         "The confidence promise is breaking"],
    ], Inches(0.55), Inches(1.6), Inches(8.9), col_widths=[1.9, 3.6, 3.0], size=12)
    box = s.shapes.add_textbox(Inches(0.7), Inches(4.75), Inches(8.6), Inches(1.6))
    bullets(box.text_frame, [
        (0, "Thresholds committed before replaying the stream — no tuning "
            "alerts until they fire.", True),
        (0, "Rolling 3-month windows: single months (6–48 sales) are too "
            "small — a 10-row window flags every feature as drifted "
            "(verified noise).", False),
    ], size=14, space_after=6)

    # ---- 20. monitoring results ------------------------------------------------------------------------------
    s = two_content_slide(prs, "2010 replay: the model ages — and says so")
    bullets(body_frame(s, 1), [
        (0, "Jan–Feb: healthy (errors below baseline, coverage on target).",
         False),
        (0, "March: RMSE spike + drift over threshold → first retrain "
            "signal.", True),
        (0, "Apr–Jul: bias marches −0.4% → −7.9% as the rebound outruns the "
            "training window; coverage sags to ~73%.", False),
        (0, "Verdict: retrain from March 2010.", True),
        (0, "Lesson: input drift says things changed; error & bias monitors "
            "prove it matters.", False),
    ], size=14)
    add_picture_fit(s, FIG / "09_monitoring_timeline.png",
                    Inches(5.15), Inches(1.5), Inches(4.55), Inches(5.5))

    # ---- 21. recommendations -------------------------------------------------------------------------------------
    s = two_content_slide(prs, "Recommendations & limitations")
    bullets(body_frame(s, 1), [
        (0, "For the platform:", True),
        (1, "Price on quality + location first; school proximity is the "
            "amenity actually priced.", False),
        (1, "Always show the range — it turns a black box into a "
            "negotiation aid.", False),
        (1, "Treat the model as perishable: drifted beyond tolerance in ~3 "
            "months; budget quarterly retraining.", False),
        (1, "Route wide-interval homes to appraisal as a premium service.",
         False),
    ], size=14)
    bullets(body_frame(s, 2), [
        (0, "Limitations:", True),
        (1, "Contextual layer is synthetic — documented and test-enforced, "
            "but real data would shift coefficients.", False),
        (1, "Prices at the 2010 Ames level; new markets need local data.",
         False),
        (1, "Cheapest decile over-predicted; segment model is future work.",
         False),
        (1, "Retraining recommended by triggers, not yet automated.", False),
    ], size=14)

    # ---- 22. closing -----------------------------------------------------------------------------------------------
    closing = prs.slides.add_slide(prs.slide_layouts[12])
    add_text(closing, Inches(4.0), Inches(2.9), Inches(5.6), Inches(1.2), [
        ("THANK YOU !", 44, True, HUST_RED)], align=PP_ALIGN.CENTER)
    links_box = add_text(closing, Inches(4.0), Inches(4.2), Inches(5.6),
                         Inches(1.0), [
        (f"Live tool: {APP_URL}", 12, False, INK),
        ("Code: github.com/ducnt3/house-price-analytics", 12, False, INK)],
        align=PP_ALIGN.CENTER)
    paras = links_box.text_frame.paragraphs
    paras[0].runs[0].hyperlink.address = f"https://{APP_URL}/"
    paras[1].runs[0].hyperlink.address = "https://github.com/ducnt3/house-price-analytics"

    # bilingual speaker notes, applied in slide-creation order
    slides = list(prs.slides)
    assert len(slides) == len(NOTES), (len(slides), len(NOTES))
    for slide, (vi, en) in zip(slides, NOTES):
        set_notes(slide, vi, en)

    out = ROOT / "reports" / "hust-house-price-slides.pptx"
    prs.save(out)
    print(f"wrote {out} ({len(prs.slides._sldIdLst)} slides, "
          f"{len(NOTES)} bilingual notes)")


if __name__ == "__main__":
    build()
